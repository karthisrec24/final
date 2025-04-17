import streamlit as st
import fitz  # PyMuPDF for PDF text extraction & image handling
import pdfplumber  # For table extraction
from PIL import Image
import io
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from langchain.llms import HuggingFaceHub
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
import pyttsx3  # Offline Text-to-Speech
import base64
import os
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import tempfile

# Streamlit UI
st.title("📖PERSONALISED TUTOR")
st.write("Upload a **PDF, DOCX, or PPTX** and ask questions based on its content!")
# Model selection
model_options = {
    "Mistral 7B": "mistralai/Mistral-7B-Instruct-v0.3",
    "Zephyr 7B": "HuggingFaceH4/zephyr-7b-beta",
    "Falcon 7B": "tiiuae/falcon-7b-instruct"
}
selected_model_name = st.sidebar.selectbox("🧠 Select Model", list(model_options.keys()))
selected_model = model_options[selected_model_name]
hf_api_key = st.sidebar.text_input("🔑 Hugging Face API Key", type="password")
st.sidebar.markdown("### 🔑 How to Generate API Key")
st.sidebar.markdown(
    "[Click here to generate your Hugging Face API Key](https://huggingface.co/settings/tokens)",
    unsafe_allow_html=True
)
if not hf_api_key:
    st.warning("Please enter your Hugging Face API key to proceed.")
    st.stop()
else:
    os.environ["HUGGINGFACEHUB_API_TOKEN"] = hf_api_key
# File uploader
uploaded_files = st.file_uploader("Upload documents (PDF, DOCX, PPTX)", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
# Initialize session state
if "faiss_index" not in st.session_state:
    st.session_state.faiss_index = None
if "text_chunks" not in st.session_state:
    st.session_state.text_chunks = None
if "qa_history" not in st.session_state:
    st.session_state.setdefault("qa_history", [])
# Load embedding model
@st.cache_resource
def load_embedding_model():
    return SentenceTransformer("all-MiniLM-L6-v2")
embedding_model = load_embedding_model()
# Extract text from PDF
def extract_text_from_pdf(pdf_bytes):
    text = ""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    for page in doc:
        text += page.get_text("text") + "\n"
    return text
# Extract tables from PDF
def extract_tables_from_pdf(pdf_bytes):
    tables = []
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages:
            table = page.extract_table()
            if table:
                # Replace None values with an empty string
                cleaned_table = [
                    [" " if cell is None else str(cell) for cell in row]  # Convert all to str
                    for row in table
                ]
                tables.append("\n".join(["\t".join(row) for row in cleaned_table]))  # Join rows
    return "\n".join(tables)
# Extract text from DOCX
def extract_text_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    return "\n".join([para.text for para in doc.paragraphs])
# Extract text from PPTX
def extract_text_from_pptx(pptx_bytes):
    presentation = Presentation(io.BytesIO(pptx_bytes))
    return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
# Store text embeddings in FAISS
@st.cache_resource
def create_faiss_index(text_data):
    text_chunks = text_data.split("\n\n")
    embeddings = np.array([embedding_model.encode(chunk) for chunk in text_chunks])
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    return index, text_chunks
# Retrieve relevant text
def retrieve_relevant_text(query, index, text_chunks, top_k=3):
    query_embedding = np.array([embedding_model.encode(query)])
    distances, indices = index.search(query_embedding, top_k)
    return [text_chunks[i] for i in indices[0]]
# Load LLM
try:
    llm = HuggingFaceHub(
        repo_id=selected_model,
        model_kwargs={"temperature": 0.5},
        callbacks=[StreamingStdOutCallbackHandler()]
    )
except Exception as e:
    st.error(f"Error loading Hugging Face model: {e}")
    st.stop()
# Define prompt
prompt = PromptTemplate(
    input_variables=["context", "question"],
    template="Answer based on the context below:\n\nContext: {context}\n\nQuestion: {question}\n\nAnswer:"
)
qa_chain = LLMChain(llm=llm, prompt=prompt)
# Generate offline audio response
def generate_audio(answer):
    engine = pyttsx3.init()
    engine.save_to_file(answer, "temp_audio.mp3")
    engine.runAndWait()
    with open("temp_audio.mp3", "rb") as f:
        audio_bytes = io.BytesIO(f.read())
    os.remove("temp_audio.mp3")
    return audio_bytes
# PDF export
def wrap_text(text, max_chars=100):
    words = text.split()
    lines, current = [], ""
    for word in words:
        if len(current) + len(word) + 1 <= max_chars:
            current += " " + word if current else word
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines

def export_qa_to_pdf(qa_history):
    temp_pdf_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name
    c = canvas.Canvas(temp_pdf_path, pagesize=letter)
    width, height = letter
    y = height - 40

    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, y, "Q&A Session History")
    y -= 30

    c.setFont("Helvetica", 12)
    for i, qa in enumerate(qa_history, 1):
        q_text = f"Q{i}: {qa['question']}"
        a_text = f"A{i}: {qa['answer'].split('Answer:')[-1].strip()}"

        for line in [q_text, a_text]:
            wrapped_lines = wrap_text(line)
            for sub_line in wrapped_lines:
                if y < 40:
                    c.showPage()
                    y = height - 40
                    c.setFont("Helvetica", 12)
                c.drawString(40, y, sub_line)
                y -= 18
            y -= 10

    c.save()
    return temp_pdf_path

# Process uploaded files
all_text_data = ""
if uploaded_files:
    for uploaded_file in uploaded_files:
        file_bytes = uploaded_file.read()
        with st.spinner(f"Processing {uploaded_file.name}... ⏳"):
            if uploaded_file.type == "application/pdf":
                extracted_text = extract_text_from_pdf(file_bytes)
                extracted_tables = extract_tables_from_pdf(file_bytes)
                all_text_data += extracted_text + "\n".join(extracted_tables)

            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                all_text_data += extract_text_from_docx(file_bytes)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                all_text_data += extract_text_from_pptx(file_bytes)
    index, text_chunks = create_faiss_index(all_text_data)
    st.session_state.faiss_index = index
    st.session_state.text_chunks = text_chunks
    st.success("✅ All documents processed! You can now ask questions.")
# Display Q&A history
if st.session_state.qa_history:
    st.write("### Previous Questions & Answers:")
    for qa in st.session_state.qa_history:
        st.markdown(f"**Q:** {qa['question']}")
        answer_text = qa['answer'].split("Answer:")[-1].strip()
        st.markdown(f"**A:** {answer_text}")
        audio_bytes = generate_audio(answer_text)
        st.audio(audio_bytes, format="audio/mp3")
    b64_audio = base64.b64encode(audio_bytes.read()).decode()
    st.markdown(f'<a href="data:audio/mp3;base64,{b64_audio}" download="answer.mp3">📥 Download Audio</a>', unsafe_allow_html=True)
    st.markdown("---")
    pdf_path = export_qa_to_pdf(st.session_state.qa_history)
    with open(pdf_path, "rb") as f:
        st.download_button(
            label="Download Q&A History as PDF",
            data=f,
            file_name="qa_history.pdf",
            mime="application/pdf"
        )

# Question input
def ask_question():
    user_question = st.session_state["user_question"]
    if not user_question.strip():
        return
    if st.session_state.faiss_index is None:
        st.warning("No document has been uploaded yet. Please upload a file first.")
        return
    relevant_chunks = retrieve_relevant_text(user_question, st.session_state.faiss_index, st.session_state.text_chunks)
    context = "\n".join(relevant_chunks)
    with st.spinner("Generating answer... ⏳"):
        answer = qa_chain.run({"context": context, "question": user_question})
    st.session_state.qa_history.append({"question": user_question, "answer": answer})
    st.session_state["user_question"] = ""
st.text_input("Ask a question about the document:", key="user_question", on_change=ask_question)